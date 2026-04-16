from __future__ import annotations

from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Dict, List

import requests
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

HEADERS = {"User-Agent": "Mozilla/5.0"}

# slide 2 exact targets from the approved template
TARGETS = {
    "page_num_top": 0,
    "header_product": 1,
    "header_index": 2,
    "header_year": 3,
    "header_brand": 4,
    "left_zone": 5,
    "left_zone_text": 6,
    "item1_image": 7,
    "item1_kind": 8,
    "item1_name": 9,
    "item1_brand_label": 10,
    "item1_brand": 11,
    "item1_model_label": 12,
    "item1_model": 13,
    "item1_size_label": 14,
    "item1_size": 15,
    "divider": 16,
    "item2_image": 17,
    "item2_kind": 18,
    "item2_name": 19,
    "item2_brand_label": 20,
    "item2_brand": 21,
    "item2_model_label": 22,
    "item2_model": 23,
    "item2_size_label": 24,
    "item2_size": 25,
    "page_num_bottom": 26,
}


def duplicate_slide(prs: Presentation, index: int):
    source = prs.slides[index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    # remove default layout placeholders
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)
    for shape in source.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def set_text(shape, text: str):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text or ""


def download_image_bytes(url: str) -> BytesIO | None:
    if not url:
        return None
    try:
        r = requests.get(url, headers=HEADERS, timeout=25)
        r.raise_for_status()
        return BytesIO(r.content)
    except Exception:
        return None


def add_image_contain(slide, image_url: str, ref_shape_index: int):
    ref = slide.shapes[ref_shape_index]
    left, top, box_w, box_h = ref.left, ref.top, ref.width, ref.height

    img_bytes = download_image_bytes(image_url)
    if not img_bytes:
        return

    try:
        img = Image.open(img_bytes)
        iw, ih = img.size
        if iw <= 0 or ih <= 0:
            return

        box_ratio = box_w / box_h
        img_ratio = iw / ih

        if img_ratio >= box_ratio:
            new_w = box_w
            new_h = int(box_w / img_ratio)
        else:
            new_h = box_h
            new_w = int(box_h * img_ratio)

        new_left = int(left + (box_w - new_w) / 2)
        new_top = int(top + (box_h - new_h) / 2)

        # remove old picture placeholder
        sp = ref._element
        sp.getparent().remove(sp)

        slide.shapes.add_picture(BytesIO(img_bytes.getvalue()), Emu(new_left), Emu(new_top), width=Emu(new_w), height=Emu(new_h))
    except Exception:
        return


def clear_item2(slide):
    keys = [
        "item2_kind", "item2_name", "item2_brand", "item2_model", "item2_size"
    ]
    for key in keys:
        set_text(slide.shapes[TARGETS[key]], "")
    # clear picture placeholder by removing it
    ref = slide.shapes[TARGETS["item2_image"]]
    sp = ref._element
    sp.getparent().remove(sp)


def fill_product_pair_slide(slide, products: List[Dict], slide_page_number: int, group_number: int, year: str):
    # page numbers / header
    set_text(slide.shapes[TARGETS["page_num_top"]], str(slide_page_number))
    set_text(slide.shapes[TARGETS["page_num_bottom"]], str(slide_page_number))
    set_text(slide.shapes[TARGETS["header_product"]], "PRODUCT")
    set_text(slide.shapes[TARGETS["header_index"]], f"{group_number:02d}")
    set_text(slide.shapes[TARGETS["header_year"]], str(year))
    set_text(slide.shapes[TARGETS["header_brand"]], "Magic Furniture")

    p1 = products[0]
    set_text(slide.shapes[TARGETS["item1_kind"]], f"{p1.get('kind', 'PRODUCT')} 01")
    set_text(slide.shapes[TARGETS["item1_name"]], p1.get("name", ""))
    set_text(slide.shapes[TARGETS["item1_brand"]], p1.get("brand", "Magic-furniture"))
    set_text(slide.shapes[TARGETS["item1_model"]], p1.get("model", p1.get("name", "")))
    set_text(slide.shapes[TARGETS["item1_size"]], p1.get("size", ""))
    add_image_contain(slide, p1.get("image_url", ""), TARGETS["item1_image"])

    if len(products) > 1:
        p2 = products[1]
        set_text(slide.shapes[TARGETS["item2_kind"]], f"{p2.get('kind', 'PRODUCT')} 02")
        set_text(slide.shapes[TARGETS["item2_name"]], p2.get("name", ""))
        set_text(slide.shapes[TARGETS["item2_brand"]], p2.get("brand", "Magic-furniture"))
        set_text(slide.shapes[TARGETS["item2_model"]], p2.get("model", p2.get("name", "")))
        set_text(slide.shapes[TARGETS["item2_size"]], p2.get("size", ""))
        add_image_contain(slide, p2.get("image_url", ""), TARGETS["item2_image"])
    else:
        clear_item2(slide)


def update_cover(slide, year: str, month: str):
    # template indices from approved cover
    set_text(slide.shapes[5], str(year))
    set_text(slide.shapes[6], f"{int(month):02d}")
    set_text(slide.shapes[13], "1")


def build_presentation(template_path: str, output_path: str, products: List[Dict], year: str, month: str):
    prs = Presentation(template_path)
    if len(prs.slides) < 2:
        raise ValueError("template.pptx must include cover + one product template slide.")

    update_cover(prs.slides[0], year, month)

    template_index = 1
    # trim to 2 slides first
    while len(prs.slides) > 2:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    pairs = [products[i:i+2] for i in range(0, len(products), 2)]
    if not pairs:
        prs.save(output_path)
        return

    # fill first template slide
    fill_product_pair_slide(prs.slides[1], pairs[0], slide_page_number=2, group_number=1, year=year)

    page_no = 3
    group_no = 2
    for pair in pairs[1:]:
        slide = duplicate_slide(prs, template_index)
        fill_product_pair_slide(slide, pair, slide_page_number=page_no, group_number=group_no, year=year)
        page_no += 1
        group_no += 1

    prs.save(output_path)
