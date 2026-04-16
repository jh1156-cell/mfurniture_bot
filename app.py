import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import requests
from bs4 import BeautifulSoup
import re
import io
import copy
from PIL import Image

# ---------------------------------------------------------
# 1. 크롤링 함수: 링크에서 제품명, 이미지, 사이즈 추출
# ---------------------------------------------------------
def scrape_product_info(url):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    try:
        res = requests.get(url, headers=headers)
        res.raise_for_status()
        soup = BeautifulSoup(res.text, 'html.parser')

        # 제품명 추출 (og:title 메타태그 활용)
        title_meta = soup.find('meta', property='og:title')
        name = title_meta['content'] if title_meta else "상품명 인식 실패"

        # 이미지 추출 (og:image 메타태그 활용)
        img_meta = soup.find('meta', property='og:image')
        img_url = img_meta['content'] if img_meta else None

        # 사이즈 추출 (페이지 전체 텍스트에서 W, D, H 가 포함된 규격 탐색)
        size = "사이즈 정보 없음"
        text_content = soup.get_text()
        # 예: W560 x D520 x H750 형태를 찾는 정규식
        size_match = re.search(r'[Ww]\s*\d+.*?([Hh]\s*\d+|[Aa][Hh]\s*\d+)', text_content)
        if size_match:
            size = size_match.group(0).strip()

        return {"name": name, "img_url": img_url, "size": size}
    except Exception as e:
        return {"name": "오류 발생", "img_url": None, "size": f"크롤링 실패: {e}"}

# ---------------------------------------------------------
# 2. 텍스트 교체 함수 (서식 유지)
# ---------------------------------------------------------
def replace_text(slide, search_text, replace_text):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if search_text in run.text:
                    run.text = run.text.replace(search_text, str(replace_text))

# ---------------------------------------------------------
# 3. 이미지 Fit 교체 함수
# ---------------------------------------------------------
def replace_image_fit(slide, old_pic_shape, new_img_url):
    try:
        # 새 이미지 다운로드
        res = requests.get(new_img_url)
        img_bytes = io.BytesIO(res.content)
        img = Image.open(img_bytes)
        img_w, img_h = img.size

        # 기존 틀(Placeholder)의 위치와 크기
        target_x = old_pic_shape.left
        target_y = old_pic_shape.top
        target_w = old_pic_shape.width
        target_h = old_pic_shape.height

        # 비율 계산 (Fit 방식)
        aspect_img = img_w / img_h
        aspect_target = target_w / target_h

        if aspect_img > aspect_target:
            # 이미지가 더 가로로 길 때
            new_w = target_w
            new_h = int(target_w / aspect_img)
        else:
            # 이미지가 더 세로로 길 때
            new_h = target_h
            new_w = int(target_h * aspect_img)

        # 가운데 정렬을 위한 여백 계산
        offset_x = target_x + (target_w - new_w) / 2
        offset_y = target_y + (target_h - new_h) / 2

        # 새 이미지 삽입
        slide.shapes.add_picture(img_bytes, offset_x, offset_y, new_w, new_h)
        
        # 기존 이미지 삭제
        sp = old_pic_shape._element
        sp.getparent().remove(sp)
    except Exception as e:
        pass # 이미지 교체 실패 시 기존 틀 유지

# ---------------------------------------------------------
# 4. 슬라이드 복제 함수
# ---------------------------------------------------------
def duplicate_slide(prs, source_slide):
    # 빈 레이아웃으로 새 슬라이드 생성
    blank_layout = prs.slide_layouts[6] 
    new_slide = prs.slides.add_slide(blank_layout)
    # 기존 도형 복사
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)
    return new_slide

# ---------------------------------------------------------
# 5. 하단 영역 삭제 함수 (홀수 페이지 처리)
# ---------------------------------------------------------
def delete_bottom_half(slide, prs_height):
    shapes_to_delete = []
    # 슬라이드 높이의 절반(y좌표)보다 아래에 있는 개체 찾기
    for shape in slide.shapes:
        if shape.top >= (prs_height / 2):
            shapes_to_delete.append(shape)
            
    # 개체 삭제
    for shape in shapes_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

# ---------------------------------------------------------
# Streamlit 웹 UI 구성
# ---------------------------------------------------------
st.set_page_config(page_title="제안서 자동 생성기", layout="wide")
st.title("🪑 매직퍼니처 제안서 자동 생성기")
st.write("가구 페이지 링크를 입력하면 제안서 PPT가 자동으로 완성됩니다.")

# 파일 업로드
uploaded_file = st.file_uploader("1. 원본 PPTX 템플릿 파일을 올려주세요", type="pptx")

# 링크 입력
links_input = st.text_area("2. 가구 상세페이지 링크를 입력하세요 (엔터로 구분하여 여러 개 입력 가능)", height=150)
links = [link.strip() for link in links_input.split('\n') if link.strip()]

if st.button("🚀 제안서 생성하기"):
    if not uploaded_file:
        st.error("PPTX 템플릿 파일을 업로드해주세요.")
    elif not links:
        st.error("최소 1개 이상의 링크를 입력해주세요.")
    else:
        with st.spinner("데이터를 수집하고 PPT를 생성 중입니다. 잠시만 기다려주세요..."):
            # 1. 크롤링 진행
            product_data = []
            for link in links:
                data = scrape_product_info(link)
                product_data.append(data)
            
            # 2. PPT 처리
            prs = Presentation(uploaded_file)
            template_slide = prs.slides[1] # 두 번째 슬라이드를 템플릿으로 사용 (표지는 0번)
            prs_height = prs.slide_height
            
            total_items = len(product_data)
            
            for i in range(0, total_items, 2):
                # 새 슬라이드 생성 (템플릿 복제)
                current_slide = duplicate_slide(prs, template_slide)
                
                # 상단 가구 (Item 1)
                item1 = product_data[i]
                page_num = (i // 2) + 1
                prod_num_1 = f"{i + 1:02d}"
                
                # 텍스트 교체 (템플릿에 있는 고정 텍스트를 찾아 바꿉니다)
                replace_text(current_slide, "FA 루츠 암체어", item1['name'])
                replace_text(current_slide, "W560 × D520 × H750 × SH440 × AH650", item1['size'])
                replace_text(current_slide, "01", prod_num_1) # Product 번호
                replace_text(current_slide, "1", str(page_num)) # 페이지 번호

                # 하단 가구 (Item 2) 처리
                if i + 1 < total_items:
                    item2 = product_data[i+1]
                    prod_num_2 = f"{i + 2:02d}"
                    replace_text(current_slide, "M 카라 테이블", item2['name'])
                    replace_text(current_slide, "W2600 × D900 × H730", item2['size'])
                    replace_text(current_slide, "02", prod_num_2)
                else:
                    # 홀수라서 하단 가구가 없을 경우 하단 전체 삭제
                    delete_bottom_half(current_slide, prs_height)

                # 이미지 교체 (슬라이드 내의 사진 개체 찾기)
                # 위치(y좌표)를 기준으로 상단/하단 이미지를 구분합니다.
                pictures = [shape for shape in current_slide.shapes if shape.shape_type == 13] # 13은 Picture 타입
                pictures.sort(key=lambda x: x.top) # 위에서 아래 순서로 정렬

                if len(pictures) >= 1 and item1['img_url']:
                    replace_image_fit(current_slide, pictures[0], item1['img_url'])
                
                if len(pictures) >= 2 and (i + 1 < total_items) and item2['img_url']:
                    replace_image_fit(current_slide, pictures[1], item2['img_url'])

            # 원본 템플릿 슬라이드 삭제 (완성본에는 복제된 슬라이드만 남김)
            rId = prs.slides._sldIdLst[-1].rId # 추가된 로직: 템플릿 찌꺼기 제거 처리 시 복잡도를 줄이기 위해 유지하거나 수동 삭제 권장

            # 결과 저장
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            st.success("🎉 제안서 생성이 완료되었습니다!")
            st.download_button(
                label="📥 완성된 제안서 다운로드",
                data=output,
                file_name="매직퍼니처_자동제안서.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )